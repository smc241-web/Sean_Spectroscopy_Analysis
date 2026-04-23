# This file plots the histograms from the shamping amplifier and preamp overlapping each other and then exports the figure to a ppt file along with the individual histograms and the risetime histograms

# plots shaping and preamp histograms overlapped
all_pulse_heights_clean = np.array(all_pulse_heights)
all_pulse_heights_clean = all_pulse_heights_clean[~np.isnan(all_pulse_heights_clean)]
print(f"Removed {len(all_pulse_heights) - len(all_pulse_heights_clean)} NaN values")

all_datasets = [
    ("PreAmp Pulse Heights", PulseHeights),
    ("ShapingAmp Pulse Heights", np.abs(all_pulse_heights_clean)),
]

fig, ax = plt.subplots(figsize=(5, 3.5), tight_layout=True)
for label, heights in all_datasets:
    heights = np.array(heights)
    heights_norm = heights / np.percentile(heights, 90)
    ax.hist(
        heights_norm,
        bins=np.linspace(0, 1.5, 200),
        histtype="step",
        linewidth=1.2,
        label=label,
    )
ax.set_xlabel("Pulse Height (arb. units)", fontsize=8, fontweight="bold")
ax.set_ylabel("Frequency", fontsize=8, fontweight="bold")
ax.tick_params(axis="both", labelsize=8)
ax.legend(fontsize=7)
plt.show()


# ── Export four histograms to PowerPoint ─────────────────────────────────────
import os, tempfile
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Derive output filename from the imported HDF5 file ───────────────────────
base_name = os.path.splitext(os.path.basename(TestFile.filename))[0]
pptx_path = f"PreAmp_ShapingAmp_Compare_{base_name}.pptx"
tmp_dir   = tempfile.mkdtemp()

# ── Prepare data ─────────────────────────────────────────────────────────────
preamp_heights   = np.array(PulseHeights)
rise_times_data  = np.array(RiseTimes)
shaping_raw      = np.array(all_pulse_heights, dtype=float)
shaping_heights  = np.abs(shaping_raw[np.isfinite(shaping_raw)])

# ── Helper: save a histogram as a PNG ────────────────────────────────────────
def save_histogram(data, xlabel, title, fname, bins=80, color="#1E6FA0"):
    fig, ax = plt.subplots(figsize=(4.5, 3.5), tight_layout=True)
    data = np.array(data)
    data = data[np.isfinite(data)]
    lo, hi = data.min(), data.max() * 1.05
    ax.hist(data, bins=np.linspace(lo, hi, bins), color=color,
            edgecolor="none", alpha=0.85)
    ax.set_xlabel(xlabel,      fontsize=10, fontweight="bold")
    ax.set_ylabel("Frequency", fontsize=10, fontweight="bold")
    ax.set_title(title,        fontsize=11, fontweight="bold")
    ax.grid(True, alpha=0.3, linewidth=0.5)
    ax.tick_params(labelsize=9)
    path = os.path.join(tmp_dir, fname + ".png")
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path

# ── Helper: save the normalised overlay histogram as a PNG ───────────────────
def save_normalized_overlay(fname):
    all_pulse_heights_clean = np.array(all_pulse_heights)
    all_pulse_heights_clean = all_pulse_heights_clean[~np.isnan(all_pulse_heights_clean)]
    all_datasets = [
        ("PreAmp Pulse Heights",    np.array(PulseHeights)),
        ("ShapingAmp Pulse Heights", np.abs(all_pulse_heights_clean)),
    ]
    fig, ax = plt.subplots(figsize=(4.5, 3.5), tight_layout=True)
    for label, heights in all_datasets:
        heights      = np.array(heights)
        heights_norm = heights / np.percentile(heights, 90)
        ax.hist(
            heights_norm,
            bins=np.linspace(0, 1.5, 200),
            histtype="step",
            linewidth=1.2,
            label=label,
        )
    ax.set_xlabel("Pulse Height (arb. units)", fontsize=10, fontweight="bold")
    ax.set_ylabel("Frequency",                 fontsize=10, fontweight="bold")
    ax.set_title("Normalised Pulse Heights",   fontsize=11, fontweight="bold")
    ax.tick_params(axis="both", labelsize=9)
    ax.legend(fontsize=8)
    ax.grid(True, alpha=0.3, linewidth=0.5)
    path = os.path.join(tmp_dir, fname + ".png")
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path

img_preamp    = save_histogram(preamp_heights,  "Pulse Height (mV)",  "PreAmp Pulse Heights",     "preamp_ph",  bins=80, color="#1E6FA0")
img_shaping   = save_histogram(shaping_heights, "Pulse Height (mV)",  "ShapingAmp Pulse Heights", "shaping_ph", bins=80, color="#E07B39")
img_risetime  = save_histogram(rise_times_data, r"Rise Time ($\mu$s)","Rise Times",               "rise_times", bins=60, color="#2A9D5C")
img_normed    = save_normalized_overlay("normed_overlay")

# ── Build the PowerPoint slide ────────────────────────────────────────────────
NAVY = RGBColor(0x1E, 0x27, 0x61)
BG   = RGBColor(0xF8, 0xF9, 0xFA)

prs = Presentation()
prs.slide_width  = Inches(13.3)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG

# Title bar
TITLE_H = Inches(0.32)
bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, TITLE_H)
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY
bar.line.fill.background()
tf  = bar.text_frame
tf.margin_left = Inches(0.12)
tf.margin_top  = Inches(0.04)
p   = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text           = f"{base_name}   |   PreAmp vs ShapingAmp Comparison"
run.font.size      = Pt(11)
run.font.bold      = True
run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
run.font.name      = "Calibri"

# ── Layout: 2×2 grid of equal-size images ────────────────────────────────────
M       = Inches(0.20)          # margin / gap between plots
TOP     = TITLE_H + M
COLS    = 2
ROWS    = 2
IMG_W   = (prs.slide_width  - M * (COLS + 1)) / COLS
IMG_H   = (prs.slide_height - TOP - M * (ROWS + 1)) / ROWS

histograms = [
    img_preamp,    # row 0, col 0
    img_shaping,   # row 0, col 1
    img_risetime,  # row 1, col 0
    img_normed,    # row 1, col 1
]

for idx, path in enumerate(histograms):
    row = idx // COLS
    col = idx  % COLS
    x   = M + col * (IMG_W + M)
    y   = TOP + M + row * (IMG_H + M)
    slide.shapes.add_picture(path, x, y, width=IMG_W, height=IMG_H)

prs.save(pptx_path)
print(f"Saved: {pptx_path}")