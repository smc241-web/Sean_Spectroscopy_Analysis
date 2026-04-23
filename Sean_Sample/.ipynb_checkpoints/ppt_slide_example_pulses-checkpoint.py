import os, io, sys, tempfile
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Output filename from HDF5 file ──────────────────────────────────────────
base_name = os.path.splitext(os.path.basename(TestFile.filename))[0]
pptx_path = base_name + '.pptx'
tmp_dir   = tempfile.mkdtemp()

def save_fig(fig, name, dpi=150):
    path = os.path.join(tmp_dir, name + '.png')
    fig.savefig(path, dpi=dpi, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

# ── 1. Preamp pulses  (SMALL) ────────────────────────────────────────────────
fig, axes = plt.subplots(1, 4, figsize=(7, 1.2), tight_layout=True)
for i in range(4):
    axes[i].set_xlabel(r"Time ($\mu$s)", fontsize=4, fontweight="bold")
    axes[i].set_ylabel("Signal (mV)",    fontsize=4, fontweight="bold")
    axes[i].set_xlim(0, 5000)
    axes[i].plot(PreAmpPulses[ExamplePulseIndices[i]], linewidth=0.1)
    axes[i].tick_params(labelsize=3.5)
fig.suptitle("Preamp Pulses", fontsize=5, fontweight="bold")
img_preamp = save_fig(fig, 'preamp')

# ── 2. Shaping amp pulses  (SMALL — same size as preamp) ─────────────────────
fig, axes = plt.subplots(1, 4, figsize=(7, 1.2), tight_layout=True)
for i in range(4):
    axes[i].set_xlabel(r"Time ($\mu$s)", fontsize=4, fontweight="bold")
    axes[i].set_ylabel("Signal (mV)",    fontsize=4, fontweight="bold")
    axes[i].set_xlim(0, Time_us[-1])
    axes[i].plot(Time_us, ShapingAmpPulses[ExamplePulseIndices[i]], linewidth=0.1)
    axes[i].tick_params(labelsize=3.5)
fig.suptitle("Shaping Amp Pulses", fontsize=5, fontweight="bold")
img_shaping = save_fig(fig, 'shaping')

# ── 3. Trigger timings  (SMALL) ──────────────────────────────────────────────
fig, ax = plt.subplots(figsize=(2.0, 1.2), tight_layout=True)
ax.plot(TriggerTimings * 1000, linewidth=0.4)
ax.set_xlabel("Trigger Number", fontsize=4, fontweight="bold")
ax.set_ylabel("Interval (ms)",  fontsize=4, fontweight="bold")
ax.tick_params(labelsize=3.5)
fig.suptitle("Trigger Timings", fontsize=5, fontweight="bold")
img_trigger = save_fig(fig, 'trigger')

# ── 4. Corrected + smoothed  (SMALL — same size as preamp/shaping) ───────────
fig, axes = plt.subplots(1, 4, figsize=(7, 1.2), tight_layout=True)
for i in range(4):
    axes[i].plot(Time_us, CorrectedPulses[i], linewidth=0.4, label='Corr')
    axes[i].plot(Time_us, SmoothedPulses[i],  linewidth=0.4, label='Smooth')
    axes[i].set_xlabel("Time (µs)", fontsize=4, fontweight="bold")
    axes[i].tick_params(labelsize=3.5)
axes[0].legend(fontsize=3, loc='upper right')
fig.suptitle("Corrected & Smoothed", fontsize=5, fontweight="bold")
img_smooth = save_fig(fig, 'smoothed')

# ── 5. Capture PulseID text ──────────────────────────────────────────────────
captured = io.StringIO()
sys.stdout = captured
AcceptedPulseList, SmoothedDerivativeList, PeakIndexList = PulseID.PeakAcceptance(
    CorrectedPulses, SmoothedPulses, PkProminence=0.1, PkHeight=0.2, Plots=False
)
sys.stdout = sys.__stdout__
pulse_id_text = captured.getvalue().strip()

# ── 6. AlphaPulsePeakCharacterisation — exactly 4 unique plots  (MEDIUM) ─────
# Close all figures first so get_fignums() only picks up new ones each call
plt.close('all')
_orig_show    = plt.show
saved_char_figs = []
seen_char_nums  = set()

def _capture_char():
    for fn in plt.get_fignums():
        if fn not in seen_char_nums:
            seen_char_nums.add(fn)
            saved_char_figs.append(plt.figure(fn))

plt.show = _capture_char

captured2 = io.StringIO()
sys.stdout = captured2
RiseTimes, PulseHeights, RejectedPulseCounter, PulseHeightErrors = \
    PulseCharacterisation.AlphaPulsePeakCharacterisation(
        AcceptedPulseList, AcceptedPulseList,
        SmoothedDerivativeList, PeakIndexList, Time_us,
        BoundaryWindow=10, BoundaryThreshold=0.1, OutsideROI=100,
        PlotNumber=4
    )
plt.show = _orig_show
sys.stdout = sys.__stdout__

img_char_paths = []
for k, f in enumerate(saved_char_figs[:4]):
    f.set_size_inches(3.2, 2.2)
    f.tight_layout()
    img_char_paths.append(save_fig(f, f'char_{k}'))

# ── 7. AlphaPlots — 3 plots  (LARGEST) ───────────────────────────────────────
plt.close('all')
saved_alpha_figs = []
seen_alpha_nums  = set()

def _capture_alpha():
    for fn in plt.get_fignums():
        if fn not in seen_alpha_nums:
            seen_alpha_nums.add(fn)
            saved_alpha_figs.append(plt.figure(fn))

plt.show = _capture_alpha
AlphaPlots(PulseHeights, RiseTimes, Time_us, base_name, nBinsPH=80, nBinsRT=60)
plt.show = _orig_show

img_alpha_paths = []
for k, f in enumerate(saved_alpha_figs):
    f.set_size_inches(3.8, 2.8)
    f.tight_layout()
    img_alpha_paths.append(save_fig(f, f'alpha_{k}'))

print("All figures saved. Building PPTX...")

# ═══════════════════════════════════════════════════════════════════════════════
# SINGLE SLIDE  13.3" × 7.5"
#
# ┌────────────────────────────────────────────────────────────────┐
# │  TITLE BAR                                                     │
# ├──────────────────────────────────────┬─────────────────────────┤
# │ [small] preamp ×4  │ trigger         │                         │
# ├────────────────────┴─────────────────┤  [medium] char plots    │
# │ [small] shaping ×4                   │        2 × 2 grid       │
# ├──────────────────────────────────────┤                         │
# │ [small] corrected+smoothed ×4        │                         │
# ├──────────────────────────────────────┴─────────────────────────┤
# │  [LARGE] AlphaPlot1 │ AlphaPlot2 │ AlphaPlot3 │ Pulse ID text  │
# └────────────────────────────────────────────────────────────────┘
# ═══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.3)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x1E, 0x27, 0x61)
BG   = RGBColor(0xF8, 0xF9, 0xFA)

slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG

# Title bar
TITLE_H = Inches(0.27)
bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, TITLE_H)
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY
bar.line.fill.background()
tf = bar.text_frame
tf.margin_left = Inches(0.1)
tf.margin_top  = Inches(0.02)
p  = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = base_name + "   |   Pulse Analysis Summary"
run.font.size      = Pt(10)
run.font.bold      = True
run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
run.font.name      = "Calibri"

# ── Master dimensions ────────────────────────────────────────────────────────
M        = Inches(0.08)
TOP      = TITLE_H + M
BOT      = prs.slide_height - M
TOTAL_H  = BOT - TOP
SW       = prs.slide_width

# Vertical split: top zone (small rows + char grid) | bottom zone (alpha plots)
ALPHA_H  = TOTAL_H * 0.42
TOP_H    = TOTAL_H - ALPHA_H - M
ALPHA_Y  = BOT - ALPHA_H
TOP_Y    = TOP

# Horizontal split in top zone: left (small rows) | right (char grid)
RIGHT_W  = Inches(4.30)
LEFT_W   = SW - RIGHT_W - 3 * M
LEFT_X   = M
RIGHT_X  = LEFT_X + LEFT_W + M

# Small rows — divide top zone into 3 equal rows
SMALL_H  = (TOP_H - 2 * M) / 3
ROW_A_Y  = TOP_Y
ROW_B_Y  = ROW_A_Y + SMALL_H + M
ROW_C_Y  = ROW_B_Y + SMALL_H + M

def add_img(path, x, y, w, h):
    slide.shapes.add_picture(path, x, y, width=w, height=h)

# ── TOP-LEFT: three small rows ───────────────────────────────────────────────
# Row A: preamp (wide) + trigger (narrow) side by side
TRIG_W   = LEFT_W * 0.22
PREAMP_W = LEFT_W - TRIG_W - M
add_img(img_preamp,  LEFT_X,                ROW_A_Y, PREAMP_W, SMALL_H)
add_img(img_trigger, LEFT_X + PREAMP_W + M, ROW_A_Y, TRIG_W,   SMALL_H)

# Row B: shaping amp
add_img(img_shaping, LEFT_X, ROW_B_Y, LEFT_W, SMALL_H)

# Row C: corrected + smoothed
add_img(img_smooth, LEFT_X, ROW_C_Y, LEFT_W, SMALL_H)

# ── TOP-RIGHT: 4 char plots in 2×2 grid ─────────────────────────────────────
cw = (RIGHT_W - M) / 2
ch = (TOP_H   - M) / 2

for k, path in enumerate(img_char_paths[:4]):
    col = k % 2
    row = k // 2
    x   = RIGHT_X + col * (cw + M)
    y   = TOP_Y   + row * (ch + M)
    add_img(path, x, y, cw, ch)

# ── BOTTOM: AlphaPlots + pulse-ID text ───────────────────────────────────────
n_alpha  = len(img_alpha_paths)
TEXT_W   = Inches(1.60)
HIST_TOT = SW - TEXT_W - 2 * M
alpha_w  = (HIST_TOT - M * (n_alpha - 1)) / max(n_alpha, 1)

for k, path in enumerate(img_alpha_paths):
    x = LEFT_X + k * (alpha_w + M)
    add_img(path, x, ALPHA_Y, alpha_w, ALPHA_H)

# Pulse-ID text
txt_x = LEFT_X + HIST_TOT + M
txb   = slide.shapes.add_textbox(txt_x, ALPHA_Y, TEXT_W, ALPHA_H)
tf    = txb.text_frame
tf.word_wrap   = True
tf.margin_left = Inches(0.05)
tf.margin_top  = Inches(0.05)

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Pulse ID\n"
run.font.size      = Pt(7)
run.font.bold      = True
run.font.color.rgb = NAVY
run.font.name      = "Calibri"

for line in pulse_id_text.split('\n'):
    p2   = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = line
    run2.font.size      = Pt(6.5)
    run2.font.name      = "Consolas"
    run2.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

prs.save(pptx_path)
print(f"Saved: {pptx_path}")
